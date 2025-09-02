# process_vendor.py
# Vendor profile processing + Problem Statement matching system
import os
import re
import json
from typing import List, Dict, Tuple
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import spacy
from keybert import KeyBERT
from rake_nltk import Rake
from sentence_transformers import SentenceTransformer, util
import numpy as np

import nltk

# Download stopwords only if not already present
nltk.download('stopwords', quiet=True)

# ------------------------------
# Step 1: Load Ontology + Transfer Matrix
# ------------------------------
def load_ontology_and_matrix(ontology_path: str, matrix_path: str) -> Tuple[Dict, Dict]:
    with open(ontology_path, 'r') as f:
        ontology = json.load(f)
    with open(matrix_path, 'r') as f:
        transfer_matrix = json.load(f)
    return ontology, transfer_matrix


ONTOLOGY_PATH = "tech_ontology.json"
TRANSFER_MATRIX_PATH = "domain_transfer.json"

ontology, transfer_matrix = load_ontology_and_matrix(ONTOLOGY_PATH, TRANSFER_MATRIX_PATH)

# Build keyword-to-domain mapping
keyword_to_domains = {}
for domain in ontology['TechnologyDomains']:
    domain_name = domain['name']
    for subdomain, details in domain['subdomains'].items():
        for category in ['tools', 'techniques', 'applications']:
            for item in details.get(category, []):
                keyword = item.lower().strip()
                if keyword not in keyword_to_domains:
                    keyword_to_domains[keyword] = []
                keyword_to_domains[keyword].append((domain_name, subdomain))

# Optional alias mapping
ALIASES = {
    "yolov5": "yolo",
    "you only look once": "yolo",
    "transformer": "bert",
    "bert-based": "bert",
}

# ------------------------------
# Step 2: Document Text Extraction
# ------------------------------
def extract_text_from_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text

def extract_text_from_ppt(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_doc(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.ppt', '.pptx']:
        return extract_text_from_ppt(file_path)
    elif ext in ['.doc', '.docx']:
        return extract_text_from_doc(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def preprocess_text(text: str) -> str:
    return re.sub(r'\s+', ' ', text).strip()

# ------------------------------
# Step 3: Keyword Extraction
# ------------------------------
nlp = spacy.load("en_core_web_sm")
kw_model = KeyBERT()
rake = Rake()

def extract_keywords(text: str, top_n: int = 20) -> List[str]:
    keybert_phrases = [kw[0].lower() for kw in kw_model.extract_keywords(
        text, keyphrase_ngram_range=(1, 3), stop_words='english', top_n=top_n
    )]

    doc = nlp(text)
    spacy_entities = [ent.text.lower() for ent in doc.ents if ent.label_ in ['ORG', 'PRODUCT', 'NORP', 'LOC']]

    rake.extract_keywords_from_text(text)
    rake_keywords = [kw.lower() for kw in rake.get_ranked_phrases()[:top_n]]

    return list(set(keybert_phrases + spacy_entities + rake_keywords))

def normalize_keywords(keywords: List[str]) -> List[str]:
    normalized = [ALIASES.get(kw.lower().strip(), kw.lower().strip()) for kw in keywords]
    return list(set(normalized))

# ------------------------------
# Step 4: Ontology Mapping
# ------------------------------
def map_keywords_to_domains(keywords: List[str]) -> Dict[str, List[str]]:
    domain_profile = {}
    for kw in keywords:
        if kw in keyword_to_domains:
            for _, subdomain in keyword_to_domains[kw]:
                domain_profile.setdefault(subdomain, []).append(kw)
    return domain_profile

# ------------------------------
# Step 5: Embedding + Semantic Similarity
# ------------------------------
semantic_model = SentenceTransformer('all-MiniLM-L6-v2')

def generate_embedding(text: str) -> np.ndarray:
    return semantic_model.encode(text)

def compute_semantic_similarity(embedding1: np.ndarray, embedding2: np.ndarray) -> float:
    return util.cos_sim(embedding1, embedding2).item()

# ------------------------------
# Step 6: Domain Transfer + Tool Overlap
# ------------------------------
subdomain_matrix = transfer_matrix['subdomain_transfer_matrix']

def calculate_domain_transfer_score(vendor_subdomains: List[str], ps_subdomains: List[str]) -> float:
    if not ps_subdomains:
        return 0.0
    scores = []
    for ps_sub in ps_subdomains:
        max_transfer = 0.0
        for ven_sub in vendor_subdomains:
            if ps_sub in subdomain_matrix and ven_sub in subdomain_matrix[ps_sub]:
                max_transfer = max(max_transfer, subdomain_matrix[ps_sub][ven_sub])
        scores.append(max_transfer)
    return np.mean(scores) if scores else 0.0

def calculate_tool_overlap(vendor_keywords: List[str], ps_keywords: List[str]) -> float:
    if not ps_keywords:
        return 0.0
    return len(set(vendor_keywords) & set(ps_keywords)) / len(ps_keywords)

# ------------------------------
# Step 7: Final Score
# ------------------------------
WEIGHTS = {'semantic': 0.4, 'domain': 0.4, 'tool': 0.2}

def calculate_final_score(semantic_score: float, domain_score: float, tool_score: float) -> float:
    return (WEIGHTS['semantic'] * semantic_score +
            WEIGHTS['domain'] * domain_score +
            WEIGHTS['tool'] * tool_score)

# ------------------------------
# Step 8: Processing Functions
# ------------------------------
def process_vendor(file_path: str) -> Dict:
    text = preprocess_text(extract_text_from_file(file_path))
    keywords = normalize_keywords(extract_keywords(text))
    domain_profile = map_keywords_to_domains(keywords)
    embedding = generate_embedding(" ".join(keywords))
    return {
        'keywords': keywords,
        'subdomains': list(domain_profile.keys()),
        'embedding': embedding,
        'domain_profile': domain_profile
    }

def process_problem_statement(ps_text: str) -> Dict:
    ps_text = preprocess_text(ps_text)
    keywords = normalize_keywords(extract_keywords(ps_text))
    domain_profile = map_keywords_to_domains(keywords)
    embedding = generate_embedding(" ".join(keywords))
    return {
        'keywords': keywords,
        'subdomains': list(domain_profile.keys()),
        'embedding': embedding,
        'domain_profile': domain_profile
    }
def process_solution(solution_data: dict) -> dict:
    """
    Process a solution JSON that is already loaded in memory as a dictionary.
    Returns structured data including keywords, subdomains, embedding, and domain profile.
    """
    # Preprocess the solution text
    ps_text = preprocess_text(solution_data.get("solution", ""))

    # Extract and normalize keywords
    keywords = normalize_keywords(extract_keywords(ps_text))

    # Map keywords to domains
    domain_profile = map_keywords_to_domains(keywords)

    # Generate embedding from keywords
    embedding = generate_embedding(" ".join(keywords))

    # Return structured data
    return {
        "title": solution_data.get("title", ""),
        "solution": ps_text,
        "keywords": keywords,
        "subdomains": list(domain_profile.keys()),
        "embedding": embedding.tolist(),
        "domain_profile": domain_profile,
    }
# ------------------------------
# Step 9: Vendor ↔ PS Matching
# ------------------------------
def match_vendor_to_ps(vendor_data: Dict, ps_data: Dict) -> Dict:
    semantic_score = compute_semantic_similarity(vendor_data['embedding'], ps_data['embedding'])
    domain_score = calculate_domain_transfer_score(vendor_data['subdomains'], ps_data['subdomains'])
    tool_score = calculate_tool_overlap(vendor_data['keywords'], ps_data['keywords'])
    final_score = calculate_final_score(semantic_score, domain_score, tool_score)
    return {
        'final_score': final_score,
        'component_scores': {
            'semantic_similarity': semantic_score,
            'domain_match': domain_score,
            'tool_overlap': tool_score
        },
        'matched_subdomains': list(set(vendor_data['subdomains']) & set(ps_data['subdomains'])),
        'matched_tools': list(set(vendor_data['keywords']) & set(ps_data['keywords'])),
        'justification': f"Matched {len(set(vendor_data['subdomains']) & set(ps_data['subdomains']))} subdomains "
                         f"and {len(set(vendor_data['keywords']) & set(ps_data['keywords']))} tools. "
                         f"Domain transfer score: {domain_score:.2f}."
    }

def match_vendors_to_ps(vendor_files: List[str], ps_text: str, vendor_names: List[str] = None) -> List[Dict]:
    ps_data = process_problem_statement(ps_text)
    results = []
    for i, file_path in enumerate(vendor_files):
        vendor_data = process_vendor(file_path)
        match_result = match_vendor_to_ps(vendor_data, ps_data)
        vendor_id = f"VENDOR_{i+1:03d}"
        vendor_name = vendor_names[i] if vendor_names else os.path.basename(file_path)
        results.append({
            'vendor_id': vendor_id,
            'vendor_name': vendor_name,
            'final_score': match_result['final_score'],
            'component_scores': match_result['component_scores'],
            'matched_domains': match_result['matched_subdomains'],
            'matched_tools': match_result['matched_tools'],
            'justification': match_result['justification']
        })
    results.sort(key=lambda x: x['final_score'], reverse=True)
    return results

def match_vendors_to_solution(vendor_files: List[str], processed_solution: Dict, vendor_names: List[str] = None) -> List[Dict]:
    results = []
    for i, file_path in enumerate(vendor_files):
        # Load processed vendor JSON
        with open(file_path, "r") as f:
            vendor_data = json.load(f)

        # Match vendor with solution
        match_result = match_vendor_to_ps(vendor_data, processed_solution)

        vendor_id = f"VENDOR_{i+1:03d}"
        vendor_name = vendor_names[i] if vendor_names else vendor_data.get("vendorName", os.path.basename(file_path))

        results.append({
            'vendor_id': vendor_id,
            'vendor_name': vendor_name,
            'final_score': match_result['final_score'],
            'component_scores': match_result['component_scores'],
            'matched_domains': match_result['matched_subdomains'],
            'matched_tools': match_result['matched_tools'],
            'justification': match_result['justification']
        })

    results.sort(key=lambda x: x['final_score'], reverse=True)
    return results